"""Shared in-memory knowledge base state and document text extraction helpers."""
from __future__ import annotations

import base64
import io
from pathlib import Path

# ── Shared state ──────────────────────────────────────────────────────────────

knowledge_base: dict = {
    "sources": [],           # list of source metadata dicts
    "extracted_texts": {},   # source_id → {"text": str, "image_data": dict | None}
}


# ── Text extraction ───────────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    parts = []
    for num, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"[Page {num}]\n{text}")
    return "\n\n".join(parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(io.BytesIO(file_bytes))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(row_text)
        if slide_texts:
            parts.append(f"[Slide {num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")


def process_image_bytes(file_bytes: bytes, mime_type: str) -> dict:
    from PIL import Image
    img = Image.open(io.BytesIO(file_bytes))
    max_size = 2048
    if max(img.size) > max_size:
        ratio = max_size / max(img.size)
        img = img.resize(
            (int(img.size[0] * ratio), int(img.size[1] * ratio)), Image.LANCZOS
        )
    buf = io.BytesIO()
    fmt = "PNG" if "png" in mime_type else "JPEG"
    img.save(buf, format=fmt)
    encoded = base64.b64encode(buf.getvalue()).decode("utf-8")
    return {
        "mime_type": mime_type,
        "data": encoded,
        "width": img.size[0],
        "height": img.size[1],
    }
